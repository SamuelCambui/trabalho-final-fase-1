from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "apresentacao_projeto_churn.pptx"

NAVY = RGBColor(14, 29, 47)
INK = RGBColor(31, 45, 61)
TEAL = RGBColor(0, 145, 150)
MINT = RGBColor(183, 235, 220)
ORANGE = RGBColor(239, 139, 72)
CREAM = RGBColor(247, 245, 239)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(105, 120, 133)
PALE_BLUE = RGBColor(224, 239, 243)
PALE_ORANGE = RGBColor(252, 232, 213)
RED = RGBColor(188, 72, 72)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def box(slide, x, y, w, h, fill, radius=False, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx


def bullet_list(slide, items, x, y, w, h, size=16, color=INK, gap=5):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return tx


def header(slide, kicker, title, page, dark=False):
    color = WHITE if dark else INK
    muted = MINT if dark else MUTED
    text(slide, kicker.upper(), 0.65, 0.38, 3.3, 0.25, 10, TEAL if not dark else MINT, True)
    text(slide, title, 0.65, 0.68, 11.6, 0.55, 28, color, True, font="Aptos Display")
    text(slide, f"FIAP Pós-Tech  |  Tech Challenge Fase 1  |  {page:02d}", 8.3, 7.12, 4.35, 0.2, 9, muted, False, align=PP_ALIGN.RIGHT)


def add_slide(kicker, title, page, dark=False):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY if dark else CREAM
    header(slide, kicker, title, page, dark)
    return slide


# 1
slide = prs.slides.add_slide(blank)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY
box(slide, 8.7, 0, 4.7, 7.5, TEAL)
box(slide, 9.45, 0.85, 2.7, 5.8, NAVY, True, MINT)
text(slide, "CHURN", 9.78, 1.35, 2.05, 0.4, 18, MINT, True, align=PP_ALIGN.CENTER)
text(slide, "26,54%", 9.53, 2.05, 2.5, 0.85, 34, WHITE, True, font="Aptos Display", align=PP_ALIGN.CENTER)
text(slide, "prevalência observada", 9.68, 3.0, 2.2, 0.3, 12, MINT, False, align=PP_ALIGN.CENTER)
for i, label in enumerate(["EDA", "MODELOS", "API"]):
    box(slide, 9.8, 4.0 + i * 0.55, 1.95, 0.34, WHITE if i == 0 else NAVY, True, WHITE)
    text(slide, label, 9.8, 4.05 + i * 0.55, 1.95, 0.18, 10, NAVY if i == 0 else WHITE, True, align=PP_ALIGN.CENTER)
text(slide, "Predição de churn\npara retenção de clientes", 0.75, 1.35, 7.3, 1.35, 34, WHITE, True, font="Aptos Display")
text(slide, "Apresentação das etapas realizadas até o momento", 0.78, 3.05, 6.5, 0.45, 18, MINT)
text(slide, "Um pipeline de Machine Learning: exploração, comparação de modelos, inferência e documentação.", 0.78, 4.12, 6.8, 0.75, 16, WHITE)
text(slide, "21/08/2026", 0.78, 6.62, 2, 0.25, 11, MINT, True)

# 2
slide = add_slide("Visão geral", "O projeto percorre o ciclo completo de ML", 2)
text(slide, "Status consolidado da entrega", 0.7, 1.55, 4.4, 0.3, 14, MUTED)
steps = [
    ("01", "Entendimento", "EDA, qualidade dos dados, métricas e baseline", TEAL),
    ("02", "Modelagem", "RF, MLP, validação cruzada e comparação", ORANGE),
    ("03", "Engenharia", "Pipeline modular, FastAPI e testes", TEAL),
    ("04", "Entrega", "Model Card, README e roteiro STAR", ORANGE),
]
for i, (number, name, desc, color) in enumerate(steps):
    x = 0.7 + i * 3.08
    box(slide, x, 2.15, 2.68, 2.5, WHITE, True, PALE_BLUE)
    text(slide, number, x + 0.18, 2.35, 0.55, 0.4, 23, color, True, font="Aptos Display")
    text(slide, name, x + 0.18, 2.95, 2.25, 0.35, 17, INK, True)
    text(slide, desc, x + 0.18, 3.45, 2.25, 0.75, 13, MUTED)
    box(slide, x + 0.18, 4.22, 0.55, 0.07, color)
text(slide, "Leitura executiva", 0.7, 5.35, 2.2, 0.3, 13, TEAL, True)
text(slide, "A fundação técnica está estruturada. As principais pendências estão na consolidação experimental do baseline, na execução do ambiente e na finalização da entrega audiovisual.", 0.7, 5.72, 11.8, 0.65, 18, INK, True)

# 3
slide = add_slide("Etapa 1", "O problema de negócio orienta a priorização de retenção", 3, dark=True)
box(slide, 0.7, 1.55, 5.6, 4.75, RGBColor(22, 46, 68), True, RGBColor(49, 83, 103))
text(slide, "PROBLEMA", 1.05, 1.9, 2.2, 0.25, 11, MINT, True)
text(slide, "Quais clientes têm maior propensão a cancelar?", 1.05, 2.28, 4.65, 0.85, 25, WHITE, True, font="Aptos Display")
text(slide, "Uso pretendido", 1.05, 3.55, 2.1, 0.25, 12, ORANGE, True)
bullet_list(slide, ["Apoiar campanhas preventivas de retenção", "Priorizar clientes para contato ou ofertas", "Manter revisão humana sobre a ação"], 1.05, 3.92, 4.65, 1.35, 15, WHITE)
box(slide, 6.75, 1.55, 5.85, 4.75, MINT, True, MINT)
text(slide, "MÉTRICAS E VALOR", 7.1, 1.9, 3.0, 0.25, 11, NAVY, True)
text(slide, "ROC-AUC + recall de churn", 7.1, 2.3, 4.7, 0.45, 23, NAVY, True, font="Aptos Display")
bullet_list(slide, ["ROC-AUC para seleção técnica", "Recall para reduzir falsos negativos", "Retenção mensal como KPI de negócio", "Custo de abordagem versus custo de perda"], 7.1, 3.15, 4.7, 1.8, 15, NAVY)
text(slide, "O ML Canvas existe, mas precisa ser alinhado aos números e ao protocolo final do projeto.", 7.1, 5.65, 4.8, 0.35, 12, NAVY, True)

# 4
slide = add_slide("Etapa 1", "A EDA revela um problema de churn moderadamente desbalanceado", 4)
text(slide, "Dataset Telco Customer Churn", 0.7, 1.48, 4.2, 0.3, 14, MUTED)
metrics = [("7.043", "clientes", TEAL), ("21", "colunas brutas", ORANGE), ("26,54%", "churn observado", TEAL)]
for i, (value, label, color) in enumerate(metrics):
    x = 0.7 + i * 2.35
    box(slide, x, 1.95, 2.03, 1.0, WHITE, True, PALE_BLUE)
    text(slide, value, x + 0.16, 2.12, 1.7, 0.35, 23, color, True, font="Aptos Display")
    text(slide, label, x + 0.16, 2.55, 1.7, 0.22, 11, MUTED)
box(slide, 0.7, 3.45, 5.8, 2.3, PALE_BLUE, True, PALE_BLUE)
text(slide, "Padrões observados", 1.0, 3.78, 2.5, 0.3, 16, INK, True)
bullet_list(slide, ["Churn: 17,98 meses de permanência média", "Não churn: 37,57 meses de permanência média", "Contrato mensal: 42,71% de churn", "Contrato de dois anos: 2,83% de churn"], 1.0, 4.25, 4.95, 1.25, 14, INK)
box(slide, 6.8, 3.45, 5.8, 2.3, PALE_ORANGE, True, PALE_ORANGE)
text(slide, "Qualidade e cuidados", 7.1, 3.78, 2.8, 0.3, 16, INK, True)
bullet_list(slide, ["TotalCharges precisa ser convertido para numérico", "customerID é removido do treinamento", "Relações são descritivas, não causais", "Notebook tem caminho absoluto e célula com erro"], 7.1, 4.25, 4.95, 1.25, 14, INK)

# 5
slide = add_slide("Etapa 2", "O pré-processamento foi encapsulado no pipeline", 5, dark=True)
text(slide, "Fluxo de treinamento", 0.75, 1.48, 3.0, 0.3, 14, MINT)
flow = [
    ("CSV", "dados brutos", TEAL),
    ("Limpeza", "TotalCharges\n+ ID", ORANGE),
    ("Split", "70 / 30\nestratificado", TEAL),
    ("Pipeline", "imputação\n+ encoding", ORANGE),
    ("Modelo", "RF / MLP", TEAL),
]
for i, (name, desc, color) in enumerate(flow):
    x = 0.75 + i * 2.45
    box(slide, x, 2.25, 1.92, 1.55, RGBColor(22, 46, 68), True, color)
    text(slide, name, x + 0.12, 2.52, 1.68, 0.3, 17, color, True, align=PP_ALIGN.CENTER)
    text(slide, desc, x + 0.12, 2.98, 1.68, 0.55, 12, WHITE, False, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        text(slide, ">", x + 2.02, 2.78, 0.35, 0.3, 23, MINT, True, align=PP_ALIGN.CENTER)
box(slide, 0.75, 4.55, 5.6, 1.15, MINT, True, MINT)
text(slide, "Validação", 1.05, 4.83, 1.35, 0.25, 14, NAVY, True)
text(slide, "5 folds com ROC-AUC e pré-processamento ajustado dentro do Pipeline, reduzindo risco de vazamento.", 2.0, 4.78, 3.9, 0.45, 14, NAVY)
box(slide, 6.75, 4.55, 5.85, 1.15, RGBColor(75, 48, 48), True, RED)
text(slide, "Escolha", 7.05, 4.83, 1.15, 0.25, 14, PALE_ORANGE, True)
text(slide, "MLP vence por apenas 0,0004 em ROC-AUC médio; a diferença não demonstra superioridade prática clara.", 7.95, 4.78, 4.25, 0.45, 14, WHITE)

# 6
slide = add_slide("Etapa 2", "A decisão do modelo depende do objetivo da operação", 6)
text(slide, "Validação cruzada", 0.7, 1.5, 2.8, 0.3, 14, MUTED)
# table
x0, y0 = 0.7, 2.0
col_widths = [3.0, 2.2, 2.1, 2.1, 2.1]
headers = ["Modelo", "ROC-AUC CV", "Accuracy", "Recall", "F1"]
rows = [
    ["MLP", "0,8481", "0,7948", "0,5223", "0,5751"],
    ["Random Forest", "0,8477", "0,7502", "0,7790", "0,6238"],
]
for j, (head, width) in enumerate(zip(headers, col_widths)):
    box(slide, x0 + sum(col_widths[:j]), y0, width, 0.55, NAVY)
    text(slide, head, x0 + sum(col_widths[:j]) + 0.08, y0 + 0.16, width - 0.16, 0.2, 11, WHITE, True, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    y = y0 + 0.58 + i * 0.62
    for j, (value, width) in enumerate(zip(row, col_widths)):
        fill = MINT if (i == 0 and j == 1) or (i == 1 and j in (3, 4)) else WHITE
        box(slide, x0 + sum(col_widths[:j]), y, width, 0.58, fill, False, CREAM)
        text(slide, value, x0 + sum(col_widths[:j]) + 0.08, y + 0.16, width - 0.16, 0.2, 14, INK, j == 0 or fill == MINT, align=PP_ALIGN.CENTER)
text(slide, "Leitura para o negócio", 0.7, 4.15, 3.0, 0.3, 16, TEAL, True)
bullet_list(slide, ["MLP: mais accuracy e precision; menos recall", "Random Forest: identifica mais clientes que podem cancelar", "A regra técnica escolhe MLP, mas a regra operacional pode preferir RF", "Próximo passo: calibrar threshold com custos reais"], 0.7, 4.58, 7.3, 1.35, 15, INK)
box(slide, 9.0, 4.15, 3.55, 1.8, PALE_ORANGE, True, PALE_ORANGE)
text(slide, "Resultado", 9.3, 4.48, 1.5, 0.25, 12, ORANGE, True)
text(slide, "MLP selecionada\npara a API", 9.3, 4.88, 2.9, 0.6, 22, INK, True, font="Aptos Display")

# 7
slide = add_slide("Etapa 3", "A API transforma o modelo em um serviço consumível", 7, dark=True)
box(slide, 0.75, 1.65, 3.2, 4.4, RGBColor(22, 46, 68), True, TEAL)
text(slide, "CLIENTE", 1.1, 2.0, 2.5, 0.25, 11, MINT, True, align=PP_ALIGN.CENTER)
text(slide, "Payload JSON\ncom dados do cliente", 1.1, 2.75, 2.5, 0.65, 20, WHITE, True, align=PP_ALIGN.CENTER)
text(slide, "Bearer token", 1.1, 4.05, 2.5, 0.3, 14, ORANGE, True, align=PP_ALIGN.CENTER)
text(slide, "Swagger / outros sistemas", 1.1, 5.15, 2.5, 0.35, 12, MINT, align=PP_ALIGN.CENTER)
text(slide, ">", 4.35, 3.55, 0.45, 0.4, 28, ORANGE, True, align=PP_ALIGN.CENTER)
box(slide, 5.0, 1.65, 3.25, 4.4, MINT, True, MINT)
text(slide, "FASTAPI", 5.38, 2.0, 2.5, 0.25, 11, NAVY, True, align=PP_ALIGN.CENTER)
text(slide, "/health", 5.45, 2.8, 2.3, 0.3, 18, NAVY, True, align=PP_ALIGN.CENTER)
text(slide, "/auth/login", 5.45, 3.45, 2.3, 0.3, 18, NAVY, True, align=PP_ALIGN.CENTER)
text(slide, "/predict", 5.45, 4.1, 2.3, 0.3, 18, NAVY, True, align=PP_ALIGN.CENTER)
text(slide, "schemas + dependências + routers", 5.35, 5.15, 2.55, 0.45, 12, NAVY, align=PP_ALIGN.CENTER)
text(slide, ">", 8.65, 3.55, 0.45, 0.4, 28, ORANGE, True, align=PP_ALIGN.CENTER)
box(slide, 9.3, 1.65, 3.25, 4.4, RGBColor(22, 46, 68), True, ORANGE)
text(slide, "RESPOSTA", 9.65, 2.0, 2.5, 0.25, 11, PALE_ORANGE, True, align=PP_ALIGN.CENTER)
text(slide, "Classe", 9.65, 2.85, 2.5, 0.3, 16, WHITE, True, align=PP_ALIGN.CENTER)
text(slide, "Yes / No", 9.65, 3.25, 2.5, 0.42, 24, ORANGE, True, font="Aptos Display", align=PP_ALIGN.CENTER)
text(slide, "Probabilidade", 9.65, 4.2, 2.5, 0.3, 16, WHITE, True, align=PP_ALIGN.CENTER)
text(slide, "0,6871", 9.65, 4.62, 2.5, 0.42, 24, MINT, True, font="Aptos Display", align=PP_ALIGN.CENTER)

# 8
slide = add_slide("Etapa 4", "A documentação registra desempenho, limites e uso responsável", 8)
items = [
    ("Model Card", "objetivo, dados, métricas, vieses, segurança e monitoramento", TEAL),
    ("README", "instalação, treinamento, endpoints, exemplo de payload e testes", ORANGE),
    ("Roteiro STAR", "estrutura de até 5 minutos para explicar a solução e demonstrar a API", TEAL),
]
for i, (name, desc, color) in enumerate(items):
    y = 1.75 + i * 1.45
    box(slide, 0.8, y, 11.75, 1.05, WHITE, True, PALE_BLUE)
    box(slide, 1.1, y + 0.28, 0.12, 0.48, color)
    text(slide, name, 1.5, y + 0.25, 2.2, 0.3, 18, INK, True, font="Aptos Display")
    text(slide, desc, 3.8, y + 0.28, 7.95, 0.4, 15, MUTED)
box(slide, 0.8, 6.15, 11.75, 0.58, PALE_ORANGE, True, PALE_ORANGE)
text(slide, "A documentação está mais madura que a execução local: os artefatos `.joblib`, dependências instaladas e vídeo ainda precisam ser finalizados.", 1.1, 6.33, 11.1, 0.2, 13, INK, True, align=PP_ALIGN.CENTER)

# 9
slide = add_slide("Próximos passos", "Seis ações fecham as lacunas da entrega", 9, dark=True)
next_steps = [
    ("01", "Alinhar o ML Canvas", "corrigir split, métricas e meta de negócio"),
    ("02", "Consolidar o baseline", "incluir Logistic Regression no relatório modular"),
    ("03", "Preparar o ambiente", "instalar dependências e executar pytest"),
    ("04", "Gerar artefatos", "treinar e validar model.joblib"),
    ("05", "Ampliar testes", "cobrir login inválido e /predict"),
    ("06", "Finalizar vídeo", "gravar, revisar e anexar o STAR"),
]
for i, (num, title, desc) in enumerate(next_steps):
    col = i % 2
    row = i // 2
    x = 0.85 + col * 6.15
    y = 1.65 + row * 1.48
    box(slide, x, y, 5.55, 1.08, RGBColor(22, 46, 68), True, RGBColor(49, 83, 103))
    text(slide, num, x + 0.28, y + 0.26, 0.55, 0.3, 17, ORANGE, True, font="Aptos Display")
    text(slide, title, x + 1.05, y + 0.2, 3.9, 0.28, 16, WHITE, True)
    text(slide, desc, x + 1.05, y + 0.58, 4.1, 0.25, 12, MINT)
text(slide, "Critério de pronto: projeto reproduzível, API demonstrável e narrativa de negócio consistente.", 0.85, 6.25, 11.5, 0.35, 17, MINT, True, align=PP_ALIGN.CENTER)

# 10
slide = add_slide("Fechamento", "A solução já tem forma de produto; falta fechar a validação", 10)
box(slide, 0.75, 1.6, 7.2, 4.55, NAVY, True, NAVY)
text(slide, "ENTREGUE ATÉ AQUI", 1.1, 1.98, 3.0, 0.25, 11, MINT, True)
bullet_list(slide, ["EDA documentada e achados de negócio", "Pipelines de pré-processamento reprodutíveis", "Random Forest e MLP comparadas por CV", "API FastAPI modular com autenticação", "Model Card, README e roteiro STAR"], 1.1, 2.48, 5.9, 2.2, 17, WHITE, gap=8)
text(slide, "Status: base técnica sólida, entrega final ainda parcial.", 1.1, 5.35, 5.9, 0.35, 16, ORANGE, True)
box(slide, 8.45, 1.6, 4.1, 4.55, MINT, True, MINT)
text(slide, "MENSAGEM PRINCIPAL", 8.8, 1.98, 3.1, 0.25, 11, NAVY, True)
text(slide, "A melhor decisão não é apenas o maior ROC-AUC.", 8.8, 2.55, 3.25, 1.05, 25, NAVY, True, font="Aptos Display")
text(slide, "É o modelo e o threshold que melhor equilibram retenção, custo de abordagem e revisão humana.", 8.8, 4.05, 3.15, 0.85, 16, NAVY)
text(slide, "Obrigado", 8.8, 5.48, 2.2, 0.3, 16, TEAL, True)

# Add speaker-note-like source footers in the slide XML is unnecessary; references are visible in the deck content.
OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(OUTPUT)
